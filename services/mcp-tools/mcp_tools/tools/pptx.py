"""`generate_pptx` -- an editable deck, built by filling a template's placeholders.

A PDF of slides is not a deck: "I need to tweak slide 4 before the meeting" is
the first thing anyone asks, and that single property is why python-pptx wins
over Marp here despite Marp being far more pleasant to generate (docs/15
section 3.1).

**Why the caps below are load-bearing.** python-pptx cannot measure text -- that
needs a rendering engine -- so it cannot know a bullet has overflowed its box.
PowerPoint's autofit is computed by PowerPoint when the file is *opened*, and
setting the autofit property in the XML resizes nothing until then. The only
reliable defence is refusing to put too much text in. Hence: cap the bullets per
slide, cap the characters per bullet, and truncate with an ellipsis rather than
shipping a slide whose text runs off the bottom.

The caps are applied by truncation, not by schema validation. docs/15 sketches
them as pydantic `max_length` constraints; a rejected call would be the right
answer against a frontier model and the wrong one here, because a local model
that gets a validation error on its seventh bullet tends to retry the whole call
rather than drop a bullet. Truncating always produces a deck.
"""

from __future__ import annotations

import contextlib
import io
from pathlib import Path

import anyio.to_thread
from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide as PptxSlide
from pptx.slide import SlideLayout
from pptx.text.text import TextFrame
from pydantic import BaseModel, Field

from mcp_tools.tools import (
    Timer,
    ToolResult,
    log_tool_call,
    ok,
    refused,
    store_artifact,
    truncate,
    unavailable,
)

MAX_SLIDES = 30
MAX_BULLETS_PER_SLIDE = 6
MAX_BULLET_CHARS = 140
MAX_TITLE_CHARS = 90
MAX_NOTES_CHARS = 700

LAYOUT_TITLE = "Title Slide"
LAYOUT_CONTENT = "Title and Content"
LAYOUT_SECTION = "Section Header"

_FURNITURE_IDX = frozenset({10, 11, 12})
"""Date, footer and slide-number placeholders. Never the body, whatever the template."""


# Three fields, matching the contract in docs/14 section 2. Layout is NOT one of
# them: it is chosen below from the shape of the content, because every argument
# is another thing a local model can get wrong and a slide with no bullets is
# unambiguously a section header.
#
# This class's docstring is emitted verbatim into `generate_pptx`'s argument
# schema, and therefore into every prompt on every turn. Keep it to one line --
# that is why the paragraph above is a comment.
class Slide(BaseModel):
    """One slide: a title, up to six bullets, and optional speaker notes."""

    title: str
    bullets: list[str] = Field(default_factory=list)
    notes: str = ""


class TemplateError(RuntimeError):
    """The template deck lacks a layout or placeholder we require.

    Raised loudly rather than writing text into whatever shape happens to be
    there, which is how a deck ends up with its body in the footer.
    """


class UnknownTemplate(ValueError):
    """No such template deck. Carries the valid names for the refusal message."""

    def __init__(self, name: str, valid: list[str]) -> None:
        super().__init__(name)
        self.name = name
        self.valid = valid


def _open_template(template: str, template_dir: Path | None) -> PresentationType:
    """Open a template deck by name.

    The name is matched against a set enumerated from disk, so no path is ever
    built from the model's string -- the same reasoning as the PDF template dict.
    With no template directory configured we fall back to python-pptx's bundled
    deck, which carries exactly the layout names we address by. That keeps the
    tool working on a fresh checkout; a real deployment mounts the team's own
    `.potx`-derived deck, which is the whole reason for choosing python-pptx.
    """
    available = sorted(p.stem for p in template_dir.glob("*.pptx")) if template_dir else []
    if template in available and template_dir is not None:
        return Presentation(str(template_dir / f"{template}.pptx"))
    if template == "default":
        return Presentation()
    raise UnknownTemplate(template, [*available, "default"])


def _layout(prs: PresentationType, name: str, fallback: str = LAYOUT_CONTENT) -> SlideLayout:
    """Resolve a layout by name.

    By name and never by index: indices differ between templates and shift the
    moment someone edits the master, and a silently wrong layout produces a deck
    that looks merely bad rather than broken.
    """
    layouts = {layout.name: layout for layout in prs.slide_layouts}
    if name in layouts:
        return layouts[name]
    if fallback in layouts:
        return layouts[fallback]
    raise TemplateError(f"template has neither a '{name}' nor a '{fallback}' layout")


def _body_text_frame(slide: PptxSlide) -> TextFrame:
    """The first real content placeholder's text frame, resolved by index.

    Resolved by `placeholder_format.idx` and never by position in the shape tree,
    and the date/footer/slide-number indices are excluded explicitly -- otherwise
    a template whose master orders shapes differently gets its body text written
    into the footer, which looks like a design problem rather than a bug.
    """
    for placeholder in slide.placeholders:
        idx = placeholder.placeholder_format.idx
        if idx != 0 and idx not in _FURNITURE_IDX and placeholder.has_text_frame:
            frame: TextFrame = placeholder.text_frame
            return frame
    raise TemplateError("layout has no body placeholder")


def _set_title(slide: PptxSlide, text: str) -> None:
    title = slide.shapes.title
    if title is None:
        raise TemplateError("layout has no title placeholder")
    title.text_frame.text = truncate(text, MAX_TITLE_CHARS)


def build_deck(
    title: str,
    subtitle: str,
    slides: list[Slide],
    *,
    template: str,
    template_dir: Path | None,
) -> bytes:
    """Fill a template deck and return the `.pptx` bytes.

    Boring by design: adding a slide from a layout copies that layout's
    placeholders, so the template does the design work and this function only
    puts text in boxes.
    """
    prs = _open_template(template, template_dir)

    opening = prs.slides.add_slide(_layout(prs, LAYOUT_TITLE))
    _set_title(opening, title)
    if subtitle:
        # A title layout without a subtitle placeholder is a legitimate template
        # choice, so this one is optional where the body placeholder is not.
        with contextlib.suppress(TemplateError):
            _body_text_frame(opening).text = truncate(subtitle, MAX_TITLE_CHARS)

    for entry in slides[:MAX_SLIDES]:
        bullets = [b for b in (bullet.strip() for bullet in entry.bullets) if b]
        layout_name = LAYOUT_CONTENT if bullets else LAYOUT_SECTION
        slide = prs.slides.add_slide(_layout(prs, layout_name))
        _set_title(slide, entry.title)

        if bullets:
            frame = _body_text_frame(slide)
            # Assigning the whole frame splits on newlines into one paragraph per
            # bullet and replaces whatever the layout put there. Cheaper than
            # clear()-then-add_paragraph(), and it touches no unannotated API.
            frame.text = "\n".join(
                truncate(bullet, MAX_BULLET_CHARS) for bullet in bullets[:MAX_BULLETS_PER_SLIDE]
            )
            for paragraph in frame.paragraphs:
                # Depth stays flat: nested bullets are how a slide quietly grows
                # past the box python-pptx cannot measure.
                paragraph.level = 0

        if entry.notes.strip():
            slide.notes_slide.notes_text_frame.text = truncate(entry.notes, MAX_NOTES_CHARS)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


async def generate_pptx(
    title: str,
    slides: list[Slide],
    template: str,
    *,
    template_dir: Path | None,
    artifact_dir: Path,
    artifact_base_url: str,
    subtitle: str = "",
) -> ToolResult:
    """Build the deck and store it. Never raises."""
    timer = Timer()
    title = title.strip()
    if not title:
        return refused("generate_pptx", "A deck title is required.")
    if not slides:
        return refused("generate_pptx", "Provide at least one slide.")

    try:
        data = await anyio.to_thread.run_sync(
            lambda: build_deck(
                title, subtitle, slides, template=template, template_dir=template_dir
            )
        )
    except UnknownTemplate as exc:
        return refused(
            "generate_pptx",
            f"Unknown template '{exc.name}'. Valid: {', '.join(exc.valid)}.",
        )
    except TemplateError as exc:
        # Our deployment is wrong, not the model's call. Say so plainly rather
        # than letting it retry a request that cannot succeed.
        log_tool_call("generate_pptx", outcome="template_error", duration_ms=timer.ms)
        return unavailable("python-pptx", f"The slide template is unusable: {exc}")

    artifact = store_artifact(
        data,
        suffix=".pptx",
        filename=title,
        directory=artifact_dir,
        base_url=artifact_base_url,
    )
    log_tool_call(
        "generate_pptx",
        outcome="ok",
        duration_ms=timer.ms,
        template=template,
        slides=min(len(slides), MAX_SLIDES),
        bytes=artifact["bytes"],
    )
    dropped = max(0, len(slides) - MAX_SLIDES)
    if dropped:
        return ok(**artifact, note=f"Deck capped at {MAX_SLIDES} slides; {dropped} were dropped.")
    return ok(**artifact)
