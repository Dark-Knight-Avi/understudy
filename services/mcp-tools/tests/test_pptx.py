"""Tests for `generate_pptx`.

python-pptx runs for real here -- it is a pure-Python library with no service
behind it, so there is no boundary to mock and nothing to gain from pretending.
The deck is written to bytes and reopened, which is as close to "opens in
PowerPoint" as a unit test can get.

The truncation tests are the ones that matter. python-pptx cannot measure text,
so a bullet that overflows produces a *valid file with a broken slide* -- there
is no error to assert on, only the cap.
"""

from __future__ import annotations

import asyncio
import io
from pathlib import Path
from typing import Any

from mcp_tools.tools.pptx import (
    MAX_BULLET_CHARS,
    MAX_BULLETS_PER_SLIDE,
    MAX_SLIDES,
    MAX_TITLE_CHARS,
    Slide,
    build_deck,
    generate_pptx,
)
from pptx import Presentation


def deck(slides: list[Slide], *, title: str = "Platform review", subtitle: str = "") -> Any:
    data = build_deck(title, subtitle, slides, template="default", template_dir=None)
    return Presentation(io.BytesIO(data))


class TestStructure:
    def test_a_title_slide_plus_one_slide_per_entry(self) -> None:
        prs = deck([Slide(title="One"), Slide(title="Two", bullets=["a"])])
        assert len(prs.slides) == 3
        assert prs.slides[0].shapes.title is not None
        assert prs.slides[0].shapes.title.text_frame.text == "Platform review"

    def test_layouts_are_chosen_by_name(self) -> None:
        prs = deck([Slide(title="Content", bullets=["a", "b"])])
        assert prs.slides[0].slide_layout.name == "Title Slide"
        assert prs.slides[1].slide_layout.name == "Title and Content"

    def test_a_slide_with_no_bullets_becomes_a_section_header(self) -> None:
        """Chosen from the shape of the content, rather than as a fourth argument."""
        prs = deck([Slide(title="Part two")])
        assert prs.slides[1].slide_layout.name == "Section Header"

    def test_bullets_become_one_paragraph_each(self) -> None:
        prs = deck([Slide(title="Findings", bullets=["alpha", "beta", "gamma"])])
        body = [p for p in prs.slides[1].placeholders if p.placeholder_format.idx == 1][0]
        assert [p.text for p in body.text_frame.paragraphs] == ["alpha", "beta", "gamma"]
        assert all(p.level == 0 for p in body.text_frame.paragraphs)

    def test_notes_are_written_to_the_notes_slide(self) -> None:
        prs = deck([Slide(title="One", bullets=["a"], notes="Mention the deadline.")])
        assert prs.slides[1].notes_slide.notes_text_frame.text == "Mention the deadline."

    def test_a_subtitle_lands_on_the_title_slide(self) -> None:
        prs = deck([Slide(title="One")], subtitle="Q3 2026")
        texts = [p.text_frame.text for p in prs.slides[0].placeholders]
        assert "Q3 2026" in texts

    def test_empty_bullets_are_dropped_rather_than_leaving_blank_lines(self) -> None:
        prs = deck([Slide(title="One", bullets=["real", "  ", ""])])
        body = [p for p in prs.slides[1].placeholders if p.placeholder_format.idx == 1][0]
        assert [p.text for p in body.text_frame.paragraphs] == ["real"]


class TestCapsAreLoadBearing:
    """python-pptx cannot measure text; PowerPoint's autofit runs only on open."""

    def test_a_400_character_bullet_is_truncated(self) -> None:
        """docs/15 acceptance test 5: truncated cleanly, no text off the slide."""
        prs = deck([Slide(title="One", bullets=["x" * 400])])
        body = [p for p in prs.slides[1].placeholders if p.placeholder_format.idx == 1][0]
        text = body.text_frame.paragraphs[0].text
        assert len(text) == MAX_BULLET_CHARS
        assert text.endswith("…")

    def test_extra_bullets_are_dropped_not_stacked(self) -> None:
        prs = deck([Slide(title="One", bullets=[f"b{i}" for i in range(20)])])
        body = [p for p in prs.slides[1].placeholders if p.placeholder_format.idx == 1][0]
        assert len(body.text_frame.paragraphs) == MAX_BULLETS_PER_SLIDE

    def test_a_long_title_is_truncated(self) -> None:
        prs = deck([Slide(title="t" * 300, bullets=["a"])])
        title = prs.slides[1].shapes.title
        assert title is not None
        assert len(title.text_frame.text) == MAX_TITLE_CHARS

    def test_the_deck_is_capped_and_says_so(self, tmp_path: Path) -> None:
        """Silently dropping slides would be worse than the cap itself."""
        slides = [Slide(title=f"S{i}", bullets=["a"]) for i in range(MAX_SLIDES + 5)]
        result = asyncio.run(
            generate_pptx(
                "Long deck",
                slides,
                "default",
                template_dir=None,
                artifact_dir=tmp_path,
                artifact_base_url="https://ai.internal/artifacts",
            )
        )
        assert result["ok"] is True
        assert "5 were dropped" in result["note"]

    def test_over_long_input_never_raises_a_validation_error(self) -> None:
        """A local model that gets rejected on its seventh bullet retries the whole call."""
        assert Slide(title="t" * 500, bullets=["b"] * 50).bullets == ["b"] * 50


class TestToolEnvelope:
    def test_stores_the_deck_and_returns_a_url(self, tmp_path: Path) -> None:
        result = asyncio.run(
            generate_pptx(
                "Platform review",
                [Slide(title="One", bullets=["a"])],
                "default",
                template_dir=None,
                artifact_dir=tmp_path,
                artifact_base_url="https://ai.internal/artifacts",
            )
        )
        assert result["ok"] is True
        assert result["filename"] == "Platform review.pptx"
        stored = tmp_path / f"{result['file_id']}.pptx"
        assert stored.exists()
        assert len(Presentation(str(stored)).slides) == 2

    def test_unknown_template_is_refused_with_the_valid_names(self, tmp_path: Path) -> None:
        result = asyncio.run(
            generate_pptx(
                "T",
                [Slide(title="One")],
                "../../etc/passwd",
                template_dir=tmp_path,
                artifact_dir=tmp_path,
                artifact_base_url="https://ai.internal/artifacts",
            )
        )
        assert result["error"] == "refused"
        assert "default" in result["detail"]

    def test_an_operator_template_is_found_by_name(self, tmp_path: Path) -> None:
        """The name is matched against a disk listing, so no path is built from it."""
        templates = tmp_path / "templates"
        templates.mkdir()
        Presentation().save(str(templates / "housestyle.pptx"))
        result = asyncio.run(
            generate_pptx(
                "T",
                [Slide(title="One", bullets=["a"])],
                "housestyle",
                template_dir=templates,
                artifact_dir=tmp_path,
                artifact_base_url="https://ai.internal/artifacts",
            )
        )
        assert result["ok"] is True

    def test_no_slides_is_refused(self, tmp_path: Path) -> None:
        result = asyncio.run(
            generate_pptx(
                "T",
                [],
                "default",
                template_dir=None,
                artifact_dir=tmp_path,
                artifact_base_url="https://ai.internal/artifacts",
            )
        )
        assert result["error"] == "refused"
        assert list(tmp_path.iterdir()) == []
